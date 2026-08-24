#!/usr/bin/env python3
"""Generate Bizcon RSVP customer/investor + technical presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

# Aurora-aligned palette
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DARK = RGBColor(0x31, 0x2E, 0x81)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x0D, 0x94, 0x88)


def set_run(run, text, size=18, bold=False, color=SLATE_700, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=SLATE_700, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # Clear default empty run content
    p.text = ""
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return box


def paragraph(tf, text, size=16, bold=False, color=SLATE_700, space_before=6, space_after=2, align=PP_ALIGN.LEFT, level=0):
    if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return p


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, label, page, total):
    # thin indigo accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.28), Inches(13.333), Inches(0.06))
    fill_shape(line, INDIGO)
    add_textbox(slide, Inches(0.5), Inches(7.35), Inches(8), Inches(0.3), label, size=10, color=SLATE_400)
    add_textbox(
        slide,
        Inches(11.2),
        Inches(7.35),
        Inches(1.6),
        Inches(0.3),
        f"{page} / {total}",
        size=10,
        color=SLATE_400,
        align=PP_ALIGN.RIGHT,
    )


def blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank)


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    fill_shape(bar, SLATE_50)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(1.15))
    fill_shape(accent, INDIGO)
    add_textbox(slide, Inches(0.5), Inches(0.28), Inches(12), Inches(0.45), title, size=28, bold=True, color=SLATE_900)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.3), subtitle, size=13, color=SLATE_500)


def bullet_block(slide, left, top, width, height, items, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        run = p.add_run()
        set_run(run, f"•  {item}", size, False, SLATE_700)
    return box


def card(slide, left, top, width, height, title, body_lines, accent=INDIGO):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, WHITE)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    fill_shape(top_bar, accent)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.35), Inches(0.35), title, size=14, bold=True, color=SLATE_900)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), width - Inches(0.35), height - Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        set_run(run, line, 12, False, SLATE_600)


def section_divider(prs, eyebrow, title, subtitle, page, total, footer):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, INDIGO_DARK)
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(0.4), eyebrow.upper(), size=12, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1), title, size=36, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.8), subtitle, size=16, color=RGBColor(0xC7, 0xD2, 0xFE))
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3), footer, size=10, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_textbox(slide, Inches(11.2), Inches(7.0), Inches(1.6), Inches(0.3), f"{page} / {total}", size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_meta = []  # for counting — we'll set total after building structure
    FOOTER = "Bizcon RSVP  ·  Confidential"

    # Pre-define total: we'll use a fixed number matching slides we create
    # Title + business (15) + divider + technical (8) + close = ~25
    TOTAL = 25

    page = 0

    def next_page():
        nonlocal page
        page += 1
        return page

    # --- 1 Title ---
    p = next_page()
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, SLATE_900)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    fill_shape(accent, INDIGO)
    add_textbox(slide, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4), "EVENT INTELLIGENCE PLATFORM", size=12, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_textbox(slide, Inches(0.9), Inches(2.5), Inches(11), Inches(0.9), "Bizcon RSVP", size=48, bold=True, color=WHITE)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(3.5),
        Inches(11),
        Inches(0.8),
        "From curated invitation to meaningful meetings —\nwith control, privacy, and proof.",
        size=18,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(0.4), "Customer & Investor Briefing  ·  Technical appendix included", size=14, color=SLATE_400)
    add_textbox(slide, Inches(0.9), Inches(6.8), Inches(8), Inches(0.3), FOOTER, size=10, color=SLATE_500)
    add_textbox(slide, Inches(11.2), Inches(6.8), Inches(1.6), Inches(0.3), f"{p} / {TOTAL}", size=10, color=SLATE_500, align=PP_ALIGN.RIGHT)

    # --- 2 Agenda ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Agenda", "What we will cover today")
    cols = [
        ("01", "The problem", "Why summits outgrow spreadsheets and RSVP forms"),
        ("02", "The product", "Lifecycle, modules, and differentiation"),
        ("03", "Trust & value", "Security, privacy, and commercial outcomes"),
        ("04", "Technical appendix", "Architecture for CTOs and technical partners"),
    ]
    for i, (num, t, d) in enumerate(cols):
        left = Inches(0.5 + i * 3.15)
        card(slide, left, Inches(1.7), Inches(3.0), Inches(4.2), f"{num}  {t}", [d], INDIGO if i < 3 else TEAL)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 3 Problem ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "The problem", "Professional summits are not open RSVP forms")
    bullet_block(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(6.0),
        Inches(5.0),
        [
            "Guest lists scattered across Excel, CRM, email, and WhatsApp",
            "Confusion between invited, registered, and on-site attendees",
            "Networking that feels random, unfair, or unmeasurable",
            "Event-day chaos: check-in queues, badge errors, room conflicts",
            "No single source of truth for sponsors, VIPs, and leadership",
        ],
        size=16,
    )
    card(
        slide,
        Inches(7.0),
        Inches(1.7),
        Inches(5.6),
        Inches(4.0),
        "The one-line truth",
        [
            "",
            "Spreadsheets scale invitations.",
            "",
            "They do not run a summit.",
            "",
            "Buyers need an operating system for",
            "curated, high-stakes events — not another",
            "public registration widget.",
        ],
        INDIGO,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 4 Who it's for ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Who it’s for", "Organisations that run serious professional events")
    audiences = [
        ("Summit organisers", ["Conference & association teams", "Corporate event ops", "Investment / sector summits"]),
        ("On-site roles", ["Registration desks", "Check-in staff", "Meeting room coordinators"]),
        ("Delegates & guests", ["VIPs & speakers", "Sponsors & exhibitors", "Investors & BD participants"]),
    ]
    for i, (t, lines) in enumerate(audiences):
        card(slide, Inches(0.5 + i * 4.2), Inches(1.6), Inches(4.0), Inches(4.3), t, [""] + lines, INDIGO)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 5 Product one-liner ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "The product in one sentence", "Multi-tenant event operations — not a ticketing widget")
    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.8),
        Inches(12),
        Inches(1.2),
        "A multi-tenant event intelligence platform that runs the full delegate lifecycle\nwith tenant isolation, roles, and measurable outcomes.",
        size=20,
        color=SLATE_700,
    )
    steps = [
        "Invite",
        "Accept",
        "Register",
        "Profile",
        "Match",
        "Meet",
        "Check in",
        "Report",
    ]
    for i, step in enumerate(steps):
        left = Inches(0.45 + i * 1.6)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.6), Inches(1.45), Inches(1.1))
        fill_shape(shape, INDIGO if i % 2 == 0 else INDIGO_DARK)
        add_textbox(slide, left, Inches(3.9), Inches(1.45), Inches(0.5), step, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        Inches(0.6),
        Inches(5.2),
        Inches(12),
        Inches(0.8),
        "Core principle: Invitation ≠ Registration ≠ Attendance.",
        size=18,
        bold=True,
        color=SLATE_900,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 6 Lifecycle ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "The delegate lifecycle", "The operating model behind every event")
    items = [
        "1. Curate invitees (import CSV / CRM)",
        "2. Send invitations by category (VIP, press, delegate…)",
        "3. Controlled registration — not open enrollment",
        "4. Attendee profiles with privacy controls",
        "5. Matchmaking & meeting requests into rooms",
        "6. Agenda, communications, and reminders",
        "7. Check-in (QR) & consistent badge printing",
        "8. Analytics, exports, and post-event proof",
    ]
    bullet_block(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.2), items, size=16)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 7 Organiser value ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Value for organisers", "One workspace. Many events. Clear control.")
    left_items = [
        "One organisation account for many summits",
        "Role-based staff access (least privilege)",
        "Invitee import and invitation workflows",
        "Categories for VIPs, speakers, media, delegates",
        "Live operational lists with search & pagination",
    ]
    right_items = [
        "Meetings, rooms, and day-of coordination",
        "Brandable badges from a single saved design",
        "Communications and reminder automations",
        "Polls and attendee engagement tools",
        "Analytics: funnel, gaps, meeting outcomes",
    ]
    bullet_block(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5), left_items, 15)
    bullet_block(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5), right_items, 15)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 8 Delegate value ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Value for delegates", "Less friction. More useful conversations.")
    cards = [
        ("Clear journey", ["Invite → accept → register", "No ambiguous status", "QR ready for event day"]),
        ("Relevant networking", ["Profiles & preferences", "Match suggestions", "Meeting requests"]),
        ("On-site confidence", ["Agenda at hand", "Polls & engagement", "Privacy respected"]),
    ]
    for i, (t, lines) in enumerate(cards):
        card(slide, Inches(0.5 + i * 4.2), Inches(1.7), Inches(4.0), Inches(4.2), t, [""] + lines, TEAL if i == 1 else INDIGO)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 9 Differentiation ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Differentiation", "Event intelligence vs. RSVP / ticketing tools")
    # table-like rows
    rows = [
        ("Typical tools", "Bizcon RSVP"),
        ("Public form / open RSVP", "Curated invitation system"),
        ("Spreadsheet guest lists", "Tenant-safe event workspace"),
        ("Generic networking apps", "Event-scoped matchmaking + rooms"),
        ("Separate badge tools", "Design + print from one source of truth"),
        ("Post-event guesswork", "Funnel, gaps, and outcomes analytics"),
    ]
    y = Inches(1.45)
    for i, (a, b) in enumerate(rows):
        bg_color = SLATE_50 if i == 0 else WHITE
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(0.72))
        fill_shape(row, bg_color)
        row.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        add_textbox(slide, Inches(0.7), y + Inches(0.18), Inches(5.5), Inches(0.4), a, size=14, bold=(i == 0), color=SLATE_900 if i == 0 else SLATE_600)
        add_textbox(slide, Inches(6.8), y + Inches(0.18), Inches(5.7), Inches(0.4), b, size=14, bold=(i == 0), color=INDIGO if i > 0 else SLATE_900)
        y += Inches(0.72)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 10 Trust ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Built for trust", "Security and privacy are product features")
    cards = [
        ("Tenant isolation", ["Organisation A never sees", "Organisation B’s events,", "delegates, or reports."]),
        ("Least privilege", ["Roles for organisers,", "registration, and check-in.", "Minimum access by default."]),
        ("Auditability", ["Invite, register, check-in,", "admin, and export actions", "are logged."]),
        ("Privacy by design", ["Collect only what’s needed.", "Attendee privacy controls.", "Sensitive fields protected."]),
    ]
    for i, (t, lines) in enumerate(cards):
        left = Inches(0.4 + (i % 4) * 3.2)
        card(slide, left, Inches(1.7), Inches(3.05), Inches(4.2), t, [""] + lines, INDIGO)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 11 Capability waves ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Capability waves", "Honest framing of what the platform covers")
    waves = [
        ("Operating core", ["Organisations & events", "Invitees & invitations", "Registration & attendees", "Check-in, staff, audit", "Basic reporting"]),
        ("Engagement layer", ["Agenda & sessions", "Communications", "Profiles & privacy", "Meetings & rooms", "Basic matchmaking"]),
        ("Intelligence & day-of", ["Advanced analytics", "AI ranking options", "Calendar sync paths", "Polls & badges", "Meeting operations"]),
    ]
    for i, (t, lines) in enumerate(waves):
        card(slide, Inches(0.5 + i * 4.2), Inches(1.6), Inches(4.0), Inches(4.6), t, [""] + lines, INDIGO if i < 2 else TEAL)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 12 Story ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Example story", "Africa Mining Summit — one operating thread")
    story = [
        "Import ~2,000 invitees from CRM or spreadsheet",
        "Invite by category (VIP, speaker, press, delegate)",
        "Delegates register and complete profiles",
        "Meetings are requested and scheduled into rooms",
        "Staff check in with QR; badges print from one design",
        "Organisers see funnel conversion and meeting outcomes",
    ]
    bullet_block(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), story, 17)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 13 Commercial ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Commercial model", "Placeholder — insert your pricing and packaging")
    cards = [
        ("Packaging ideas", ["Per organisation", "Per event", "Per attendee tiers", "Pilot / annual contracts"]),
        ("Add-on modules", ["AI matchmaking", "CRM sync", "Badge desk", "SMS / WhatsApp (later)"]),
        ("Buyer personas", ["Ops / event director", "Commercial / sponsors", "IT / security reviewer"]),
    ]
    for i, (t, lines) in enumerate(cards):
        card(slide, Inches(0.5 + i * 4.2), Inches(1.7), Inches(4.0), Inches(4.2), t, [""] + lines, INDIGO)
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4), "Replace this slide with confirmed ACV, trial terms, and packaging before customer meetings.", size=12, color=SLATE_500)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 14 Why now ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Why now", "Market timing meets product readiness")
    bullet_block(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Professional summits are competitive again — networking ROI is a buying criterion",
            "Buyers expect measurable outcomes, not just attendance counts",
            "Privacy and tenant isolation are table stakes for enterprise and government-adjacent events",
            "AI matchmaking only works when invitation, profile, and meeting ops are already solid",
            "Bizcon RSVP is built foundation-first: secure lifecycle first, intelligence on top",
        ],
        16,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 15 Ask ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "The ask / next step", "Clear outcomes for every audience")
    card(
        slide,
        Inches(0.5),
        Inches(1.6),
        Inches(6.0),
        Inches(4.5),
        "For customers",
        [
            "",
            "• Pilot on the next summit",
            "• Structured onboarding (≈2 weeks)",
            "• Success metrics:",
            "   – invitation completion",
            "   – on-site check-in rate",
            "   – meetings scheduled & held",
            "",
            "Start with one flagship event.",
        ],
        INDIGO,
    )
    card(
        slide,
        Inches(6.8),
        Inches(1.6),
        Inches(6.0),
        Inches(4.5),
        "For investors",
        [
            "",
            "• Capital to deepen GTM",
            "• Enterprise sales motion",
            "• AI matchmaking quality",
            "• Mobile / offline day-of",
            "",
            "Insert: raise amount, use of funds,",
            "traction, and runway targets.",
        ],
        TEAL,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- 16 Close business ---
    p = next_page()
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, SLATE_900)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    fill_shape(accent, INDIGO)
    add_textbox(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.5), "Curated invitation → controlled registration\n→ meaningful meetings → measurable outcomes.", size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.9), Inches(4.6), Inches(11), Inches(0.5), "Bizcon RSVP — Event intelligence for professional summits.", size=16, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_textbox(slide, Inches(0.9), Inches(6.8), Inches(8), Inches(0.3), FOOTER, size=10, color=SLATE_500)
    add_textbox(slide, Inches(11.2), Inches(6.8), Inches(1.6), Inches(0.3), f"{p} / {TOTAL}", size=10, color=SLATE_500, align=PP_ALIGN.RIGHT)

    # --- Divider technical ---
    p = next_page()
    section_divider(
        prs,
        "Appendix",
        "Technical overview",
        "For CTOs, technical partners, and diligence conversations.",
        p,
        TOTAL,
        FOOTER,
    )

    # --- T1 Architecture ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Architecture snapshot", "Modern multi-tenant SaaS stack")
    items = [
        "Multi-tenant model: Organisation → many Events",
        "Next.js application with server-side authorization",
        "PostgreSQL via Prisma",
        "Authentication with Clerk; permissions on sensitive actions",
        "Background jobs for email and pipelines (Inngest)",
        "Integrations: HubSpot / Salesforce import; Google / Microsoft calendar paths",
    ]
    bullet_block(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), items, 16)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- T2 Authz ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Authorization model", "The strongest technical story for enterprise buyers")
    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.7),
        Inches(12),
        Inches(0.6),
        "identity  →  organisation  →  event  →  role  →  resource  →  field",
        size=20,
        bold=True,
        color=INDIGO,
        align=PP_ALIGN.CENTER,
    )
    bullet_block(
        slide,
        Inches(0.5),
        Inches(2.6),
        Inches(12),
        Inches(4),
        [
            "Frontend is never the security boundary",
            "Tenant scope always derived from the authenticated session (organisationId)",
            "Least privilege for organisers, registration, and check-in staff",
            "Check-in responses expose only the fields staff need",
            "Tokens hashed; sensitive secrets never logged",
        ],
        16,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- T3 Data model ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Lifecycle data model", "Separation is the product moat")
    entities = [
        ("Invitees", "Contacts for an event — not yet registered"),
        ("Invitations", "Category-scoped invite state & delivery"),
        ("Attendees", "Registered / confirmed participants"),
        ("Meetings", "Requests, rooms, schedules, outcomes"),
        ("Check-in / Badges", "On-site attendance & print design"),
    ]
    for i, (t, d) in enumerate(entities):
        y = Inches(1.5 + i * 0.95)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(0.85))
        fill_shape(shape, SLATE_50)
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        add_textbox(slide, Inches(0.8), y + Inches(0.22), Inches(3), Inches(0.4), t, size=16, bold=True, color=SLATE_900)
        add_textbox(slide, Inches(4.0), y + Inches(0.22), Inches(8.5), Inches(0.4), d, size=15, color=SLATE_600)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- T4 Matchmaking ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Matchmaking & meetings", "Rules + optional AI, constrained by real ops")
    bullet_block(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Structured scoring from profile preferences (looking for / offering / industries…)",
            "Optional AI ranking layered on top of structured scores",
            "Always event-scoped — no cross-event leakage by default",
            "Privacy-aware inputs; organiser can disable AI features per event",
            "Meetings constrained by rooms, conflicts, and agenda commitments",
            "Ops views: today list, room board, moderation, demand heatmaps",
        ],
        16,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- T5 Event day ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Event-day reliability", "Designed for desks that can’t fail")
    cards = [
        ("QR check-in", ["Staff see only needed fields", "Rate-limited scans", "Audit logged"]),
        ("Badges", ["Config stored in database", "Same design every printer", "Live preview & fonts/colors"]),
        ("Ops boards", ["Room grid & today view", "Conflict awareness", "Moderation tools"]),
    ]
    for i, (t, lines) in enumerate(cards):
        card(slide, Inches(0.5 + i * 4.2), Inches(1.7), Inches(4.0), Inches(4.2), t, [""] + lines, INDIGO)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- T6 Analytics ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Analytics", "Proof for organisers, sponsors, and leadership")
    items = [
        "Lifecycle funnel: invited → accepted → registered → checked in",
        "Matchmaking ROI: profiles → requests → scheduled → completed",
        "Room utilization timeline",
        "Category mix and meeting outcomes trends",
        "Gap finder: completed profiles with zero meetings",
        "CSV export for selected slices",
        "Operational tables use shared searchable / paginated DataTable UX",
    ]
    bullet_block(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), items, 16)
    add_footer(slide, FOOTER, p, TOTAL)

    # --- T7 Extensibility ---
    p = next_page()
    slide = blank_slide(prs)
    title_bar(slide, "Extensibility & quality", "Modular product surface")
    bullet_block(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Domain modules under a clear server-side structure (events, meetings, badges, polls…)",
            "CRM import paths (HubSpot / Salesforce) for invitee acquisition",
            "Communications automations and meeting reminders",
            "Shared Aurora design system for a consistent organiser UX",
            "Targeted tests around scoring, analytics, badges, and critical flows",
            "Built for 5,000–10,000 person professional summits without manual intervention in normal ops",
        ],
        16,
    )
    add_footer(slide, FOOTER, p, TOTAL)

    # --- Final close ---
    p = next_page()
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, INDIGO_DARK)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5), "Thank you", size=40, bold=True, color=WHITE)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(3.3),
        Inches(11.5),
        Inches(1.2),
        "Let’s pick a pilot event — or continue diligence with a technical deep dive.",
        size=18,
        color=RGBColor(0xC7, 0xD2, 0xFE),
    )
    add_textbox(slide, Inches(0.9), Inches(5.0), Inches(11), Inches(0.4), "Bizcon RSVP  ·  Event intelligence for professional summits", size=14, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_textbox(slide, Inches(0.9), Inches(6.8), Inches(8), Inches(0.3), FOOTER, size=10, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_textbox(slide, Inches(11.2), Inches(6.8), Inches(1.6), Inches(0.3), f"{p} / {TOTAL}", size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

    assert page == TOTAL, f"Page count mismatch: {page} vs {TOTAL}"

    out = Path(__file__).resolve().parents[1] / "docs" / "Bizcon-RSVP-Customer-Investor-Briefing.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")
    return out


if __name__ == "__main__":
    build()
